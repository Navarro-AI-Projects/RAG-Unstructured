import PyPDF2
from os import listdir
from os.path import isfile, join,isdir
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_qdrant import Qdrant
import sys
from langchain_text_splitters import TokenTextSplitter
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
import docx
from dotenv import load_dotenv
import os

load_dotenv()
def get_files(dir):
    file_list = []
    for f in listdir(dir):
        if isfile(join(dir,f)):
            file_list.append(join(dir,f))
        elif isdir(join(dir,f)):
            file_list= file_list + get_files(join(dir,f))
    return file_list

def getTextFromWord(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def getTextFromPPTX(filename):
    prs = Presentation(filename)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            fullText.append(shape.text)
    return '\n'.join(fullText)


class Document:
    def __init__(self, page_content, metadata):
        self.page_content = page_content
        self.metadata = metadata
def main_indexing(mypath):
    model_name = "models/embedding-001"
    hf = GoogleGenerativeAIEmbeddings(
        model=model_name
    )
    onlyfiles = get_files(mypath)
    file_content = ""
    qdrant = None

    local_qdrant_path = os.path.join(os.path.expanduser("~"), "local_qdrant")

    for file in onlyfiles:
        file_content = ""
        if file.endswith(".pdf"):
            print("indexing "+file)
            reader = PyPDF2.PdfReader(file)
            for i in range(0,len(reader.pages)):
                file_content = file_content + " "+reader.pages[i].extract_text()
        elif file.endswith(".txt"):
            print("indexing " + file)
            f = open(file,'r')
            file_content = f.read()
            f.close()
        elif file.endswith(".docx"):
            print("indexing " + file)
            file_content = getTextFromWord(file)
        elif file.endswith(".pptx"):
            print("indexing " + file)
            file_content = getTextFromPPTX(file)
        else:
            continue
        text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
        texts = text_splitter.split_text(file_content)
        metadata = {"path":file}
        documents = [Document(text,metadata) for text in texts]
        collection_name = "MyCollection"
        if qdrant is None:
            qdrant = Qdrant.from_documents(
                                    documents,
                                    hf,
                                    path= local_qdrant_path,
                                    collection_name=collection_name
                                    )
        else:
            qdrant.add_documents(documents)
    len(texts)
    print(onlyfiles)
    print("Finished indexing!")
